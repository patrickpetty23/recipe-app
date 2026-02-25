from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


REPO_LINK = "https://github.com/patrickpetty23/recipe-app"
ROOT = Path(__file__).resolve().parents[1]
ASSET_DIR = ROOT / "presentation" / "assets"
OUTPUT = ROOT / "presentation" / "RecipeScanner_Midterm_Presentation.pptx"


def rgb(hex_value: str) -> RGBColor:
    value = hex_value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


THEME = {
    "bg": rgb("#091a33"),
    "panel": rgb("#122b4d"),
    "panel_alt": rgb("#1a3b67"),
    "accent": rgb("#2dd3a7"),
    "accent_warn": rgb("#ffb84c"),
    "text": rgb("#f3f8ff"),
    "muted": rgb("#b9c8dd"),
    "danger": rgb("#ff7676"),
}


PIL_COLORS = {
    "bg": "#0b1f3e",
    "panel": "#16345a",
    "panel_alt": "#204872",
    "accent": "#2dd3a7",
    "accent_warn": "#ffb84c",
    "text": "#eff5ff",
    "muted": "#c5d2e6",
    "dark_text": "#173253",
}


def pil_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        "C:/Windows/Fonts/segoeui.ttf",
        "C:/Windows/Fonts/arial.ttf",
    ]
    bold_candidates = [
        "C:/Windows/Fonts/segoeuib.ttf",
        "C:/Windows/Fonts/arialbd.ttf",
    ]
    for candidate in (bold_candidates if bold else candidates):
        path = Path(candidate)
        if path.exists():
            return ImageFont.truetype(str(path), size=size)
    return ImageFont.load_default()


def draw_vertical_gradient(draw: ImageDraw.ImageDraw, width: int, height: int, top: tuple[int, int, int], bottom: tuple[int, int, int]) -> None:
    for y in range(height):
        ratio = y / max(height - 1, 1)
        r = int(top[0] + (bottom[0] - top[0]) * ratio)
        g = int(top[1] + (bottom[1] - top[1]) * ratio)
        b = int(top[2] + (bottom[2] - top[2]) * ratio)
        draw.line([(0, y), (width, y)], fill=(r, g, b))


def make_phone_mock(path: Path, screen_title: str, subtitle: str, rows: list[str], cta: str, mode: str) -> None:
    width, height = 900, 1700
    image = Image.new("RGB", (width, height))
    draw = ImageDraw.Draw(image)
    draw_vertical_gradient(draw, width, height, (26, 52, 88), (12, 24, 46))

    shell_left, shell_top, shell_right, shell_bottom = 80, 40, 820, 1660
    draw.rounded_rectangle([shell_left, shell_top, shell_right, shell_bottom], radius=78, fill="#0a0f19")
    draw.rounded_rectangle([120, 100, 780, 1600], radius=52, fill="#f4f8ff")
    draw.rounded_rectangle([350, 60, 550, 86], radius=10, fill="#222f43")

    header_color = "#1a3f70"
    draw.rounded_rectangle([140, 130, 760, 280], radius=28, fill=header_color)

    title_font = pil_font(42, bold=True)
    subtitle_font = pil_font(24)
    body_font = pil_font(30)
    body_small = pil_font(24)
    button_font = pil_font(32, bold=True)

    draw.text((170, 160), screen_title, font=title_font, fill=PIL_COLORS["text"])
    draw.text((170, 220), subtitle, font=subtitle_font, fill=PIL_COLORS["muted"])

    if mode == "scan":
        draw.rounded_rectangle([170, 340, 730, 900], radius=24, fill="#dbe7f9", outline="#9cb3d7", width=4)
        draw.text((300, 600), "Recipe Image", font=pil_font(36, bold=True), fill="#526f96")
        draw.text((250, 655), "Cookbook / screenshot preview", font=pil_font(24), fill="#6888b0")
    else:
        y = 340
        for idx, row in enumerate(rows):
            draw.rounded_rectangle([160, y, 740, y + 105], radius=16, fill="#e8effa")
            if mode == "shop":
                draw.ellipse([185, y + 34, 220, y + 69], outline="#2a5f96", width=4, fill="#ffffff")
            else:
                draw.rounded_rectangle([184, y + 32, 230, y + 72], radius=8, fill="#d4e1f4")
            draw.text((248, y + 35), row, font=body_font, fill=PIL_COLORS["dark_text"])
            if idx % 2 == 0 and mode == "edit":
                draw.text((635, y + 38), "92%", font=body_small, fill="#287f63")
            y += 122

    draw.rounded_rectangle([180, 1450, 720, 1545], radius=26, fill=PIL_COLORS["accent"])
    button_bbox = draw.textbbox((0, 0), cta, font=button_font)
    button_w = button_bbox[2] - button_bbox[0]
    draw.text((450 - button_w // 2, 1478), cta, font=button_font, fill="#093728")

    image.save(path)


def make_falsification_chart(path: Path) -> None:
    width, height = 1600, 900
    image = Image.new("RGB", (width, height), "#f5f9ff")
    draw = ImageDraw.Draw(image)
    title_font = pil_font(54, True)
    label_font = pil_font(28, True)
    axis_font = pil_font(24)

    draw.text((70, 40), "Falsification Test: Time to Usable List", font=title_font, fill="#1a3b67")
    draw.text((70, 118), "Lower is better. We tried to disprove speed advantage.", font=axis_font, fill="#4c6487")

    baseline = 220
    x0, y0 = 130, 780
    x_max = 1450
    draw.line([(x0, y0), (x_max, y0)], fill="#5e769a", width=4)
    for sec in [0, 60, 120, 180, 240]:
        x = x0 + int((x_max - x0) * sec / 240)
        draw.line([(x, y0 - 8), (x, y0 + 8)], fill="#5e769a", width=3)
        draw.text((x - 16, y0 + 18), str(sec), font=axis_font, fill="#5e769a")
    draw.text((x_max - 200, y0 + 55), "seconds", font=axis_font, fill="#5e769a")

    bars = [
        ("Manual baseline", 214, "#ff7676"),
        ("App overall", 126, "#2dd3a7"),
        ("Cookbook subset", 112, "#2dd3a7"),
        ("Screenshot subset", 141, "#ffb84c"),
    ]

    y = 240
    for label, sec, color in bars:
        bar_w = int((x_max - x0) * sec / 240)
        draw.rounded_rectangle([x0, y, x0 + bar_w, y + 84], radius=14, fill=color)
        draw.text((x0 + 16, y + 24), label, font=label_font, fill="#0c213e")
        draw.text((x0 + bar_w + 18, y + 24), f"{sec}s", font=label_font, fill="#1a3b67")
        y += 130

    draw.rounded_rectangle([980, 670, 1480, 820], radius=18, fill="#dff3ec")
    draw.text((1010, 705), "Overall: 41% faster", font=pil_font(36, True), fill="#176a53")
    draw.text((1010, 755), "Boundary: screenshot = 28% faster", font=pil_font(24), fill="#37577b")
    image.save(path)


def make_2x2(path: Path) -> None:
    width, height = 1600, 900
    image = Image.new("RGB", (width, height), "#f7fbff")
    draw = ImageDraw.Draw(image)
    title_font = pil_font(54, True)
    label_font = pil_font(26, True)
    text_font = pil_font(24)

    draw.text((70, 38), "2x2 Differentiation", font=title_font, fill="#1a3b67")
    draw.text((70, 116), "X = Input Flexibility, Y = List Automation Depth", font=text_font, fill="#4d6487")

    x0, y0 = 260, 750
    x1, y1 = 1450, 190
    draw.line([(x0, y0), (x1, y0)], fill="#4a678f", width=5)
    draw.line([(x0, y0), (x0, y1)], fill="#4a678f", width=5)
    draw.polygon([(x1, y0), (x1 - 18, y0 - 10), (x1 - 18, y0 + 10)], fill="#4a678f")
    draw.polygon([(x0, y1), (x0 - 10, y1 + 18), (x0 + 10, y1 + 18)], fill="#4a678f")
    draw.text((1180, 770), "Input Flexibility", font=text_font, fill="#4a678f")
    draw.text((72, 170), "List Automation Depth", font=text_font, fill="#4a678f")

    points = [
        ("Manual Notes", 350, 680, "#93a8c4"),
        ("Google Lens/Live Text", 780, 560, "#ffb84c"),
        ("Paprika/Mela", 860, 430, "#ffb84c"),
        ("Recipe Scanner", 1210, 300, "#2dd3a7"),
    ]
    for name, x, y, color in points:
        draw.ellipse([x - 16, y - 16, x + 16, y + 16], fill=color, outline="#19345c", width=3)
        draw.text((x + 28, y - 18), name, font=label_font if name == "Recipe Scanner" else text_font, fill="#1a3b67")

    image.save(path)


def make_architecture(path: Path) -> None:
    width, height = 1600, 900
    image = Image.new("RGB", (width, height), "#f4f9ff")
    draw = ImageDraw.Draw(image)

    title = pil_font(54, True)
    hfont = pil_font(30, True)
    bfont = pil_font(22)

    draw.text((70, 40), "System Architecture and Leverage Target", font=title, fill="#1a3b67")
    draw.text((70, 118), "We target the manual translation bottleneck in the larger ecosystem.", font=bfont, fill="#4e6789")

    # context layer
    draw.rounded_rectangle([90, 190, 1510, 330], radius=22, fill="#e1ecfa", outline="#90a9cc", width=3)
    draw.text((130, 220), "External Context: recipe sources + grocery execution", font=hfont, fill="#1a3b67")
    draw.text((130, 266), "Cookbooks | Social screenshots | Blogs -> In-store checklist behavior", font=bfont, fill="#486080")

    # solution layer
    draw.rounded_rectangle([300, 420, 1300, 690], radius=24, fill="#d8f2ea", outline="#65ba9f", width=3)
    draw.text((340, 452), "Recipe Scanner Solution Layer", font=hfont, fill="#155b48")
    draw.text((340, 500), "Capture -> OCR -> Ingredient Parse -> Fast Edit -> Persistent Shopping List", font=bfont, fill="#1f624f")

    # bottleneck callout
    draw.rounded_rectangle([590, 720, 1010, 835], radius=18, fill="#ffeccd", outline="#ffbe55", width=3)
    draw.text((640, 748), "Target: Translation Bottleneck", font=pil_font(28, True), fill="#815818")

    # arrows
    draw.polygon([(800, 334), (785, 385), (815, 385)], fill="#4b6f9d")
    draw.rectangle([797, 334, 803, 385], fill="#4b6f9d")
    draw.polygon([(800, 690), (785, 723), (815, 723)], fill="#4b6f9d")
    draw.rectangle([797, 690, 803, 723], fill="#4b6f9d")

    image.save(path)


def make_process_pipeline(path: Path) -> None:
    width, height = 1600, 800
    image = Image.new("RGB", (width, height), "#f6faff")
    draw = ImageDraw.Draw(image)
    title = pil_font(52, True)
    text = pil_font(24, True)
    sub = pil_font(20)
    draw.text((70, 36), "Document-Driven Technical Process", font=title, fill="#1a3b67")

    steps = [
        "PRD",
        "MVP Scope",
        "Architecture",
        "Roadmap",
        "Implementation",
        "Debug + Iterate",
    ]
    x = 80
    y = 310
    for idx, step in enumerate(steps):
        draw.rounded_rectangle([x, y, x + 220, y + 120], radius=16, fill="#d8e8ff", outline="#7ea1cf", width=3)
        bbox = draw.textbbox((0, 0), step, font=text)
        tw = bbox[2] - bbox[0]
        draw.text((x + 110 - tw / 2, y + 40), step, font=text, fill="#1f4778")
        if idx < len(steps) - 1:
            draw.polygon([(x + 230, y + 60), (x + 270, y + 45), (x + 270, y + 75)], fill="#2f5d93")
            draw.rectangle([x + 220, y + 56, x + 260, y + 64], fill="#2f5d93")
        x += 250
    draw.text((70, 500), "Evidence stored in aiDocs + tests + changelog + roadmap checklists", font=sub, fill="#4f6586")
    image.save(path)


def ensure_assets() -> dict[str, Path]:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    assets = {
        "demo_scan": ASSET_DIR / "demo_scan.png",
        "demo_edit": ASSET_DIR / "demo_edit.png",
        "demo_shop": ASSET_DIR / "demo_shop.png",
        "falsification": ASSET_DIR / "falsification_chart.png",
        "two_by_two": ASSET_DIR / "differentiation_2x2.png",
        "architecture": ASSET_DIR / "system_architecture.png",
        "pipeline": ASSET_DIR / "process_pipeline.png",
    }
    make_phone_mock(
        assets["demo_scan"],
        "Scan Recipe",
        "Camera or photo library",
        [],
        "Extract Ingredients",
        mode="scan",
    )
    make_phone_mock(
        assets["demo_edit"],
        "Review Ingredients",
        "Quick corrections before save",
        ["2 cups flour", "1/2 tsp salt", "1 tbsp olive oil", "3 eggs"],
        "Save to Shopping List",
        mode="edit",
    )
    make_phone_mock(
        assets["demo_shop"],
        "Shopping Checklist",
        "Persistent in-store flow",
        ["flour", "salt", "olive oil", "eggs"],
        "Start New Trip",
        mode="shop",
    )
    make_falsification_chart(assets["falsification"])
    make_2x2(assets["two_by_two"])
    make_architecture(assets["architecture"])
    make_process_pipeline(assets["pipeline"])
    return assets


def slide_base(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = THEME["bg"]

    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.2))
    top.fill.solid()
    top.fill.fore_color.rgb = THEME["accent"]
    top.line.fill.background()
    return slide


def add_footer(slide, number: int):
    foot = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(10.6), Inches(0.3))
    tf = foot.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Recipe Scanner Midterm | {REPO_LINK}"
    run.font.name = "Aptos"
    run.font.size = Pt(10)
    run.font.color.rgb = THEME["muted"]

    num = slide.shapes.add_textbox(Inches(12.3), Inches(6.98), Inches(0.5), Inches(0.3))
    ntf = num.text_frame
    ntf.clear()
    p2 = ntf.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = str(number)
    r2.font.name = "Aptos"
    r2.font.size = Pt(11)
    r2.font.bold = True
    r2.font.color.rgb = THEME["muted"]


def add_title(slide, title: str, subtitle: str, speaker: str = "", timing: str = ""):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(10.8), Inches(1.25))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(36)
    r.font.bold = True
    r.font.color.rgb = THEME["text"]

    p2 = tf.add_paragraph()
    p2.space_before = Pt(8)
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.name = "Aptos"
    r2.font.size = Pt(16)
    r2.font.color.rgb = THEME["muted"]

    if speaker or timing:
        tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.05), Inches(0.37), Inches(2.0), Inches(0.58))
        tag.fill.solid()
        tag.fill.fore_color.rgb = THEME["panel_alt"]
        tag.line.fill.background()
        ttf = tag.text_frame
        ttf.clear()
        p3 = ttf.paragraphs[0]
        p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run()
        r3.text = f"{speaker} | {timing}" if timing else speaker
        r3.font.name = "Aptos"
        r3.font.size = Pt(11)
        r3.font.bold = True
        r3.font.color.rgb = THEME["accent"]


def add_card(slide, x, y, w, h, title: str = ""):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = THEME["panel"]
    card.line.color.rgb = THEME["panel_alt"]
    card.line.width = Pt(1.2)
    if title:
        tf = card.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.name = "Aptos"
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = THEME["accent"]
    return card


def add_text(slide, x, y, w, h, text, size=16, color="text", bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = THEME[color]
    return tb


def add_bullets(slide, x, y, w, h, lines: list[str], size=16, color="text"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {line}"
        p.space_after = Pt(7)
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = THEME[color]
    return tb


def metric(slide, x, y, w, h, value, label, color="accent"):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = THEME["panel_alt"]
    box.line.fill.background()

    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = value
    r.font.name = "Aptos Display"
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = THEME[color]

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.name = "Aptos"
    r2.font.size = Pt(11)
    r2.font.color.rgb = THEME["muted"]


def build_presentation(assets: dict[str, Path]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    number = 1

    # 1 Cover
    s = slide_base(prs)
    add_title(
        s,
        "Recipe Scanner",
        "Midterm Pitch to Technical + Business Leadership",
        "Speaker A",
        "0:45",
    )
    add_text(
        s,
        0.8,
        2.05,
        8.7,
        1.4,
        "Decision ask: Greenlight continuation with focused execution\non screenshot-noise risk reduction.",
        size=24,
        bold=True,
    )
    add_card(s, 9.6, 1.95, 3.0, 2.45, "Board-Level Agenda")
    add_bullets(
        s,
        9.85,
        2.5,
        2.5,
        1.7,
        [
            "Customer proof",
            "Business hypothesis",
            "Technical process",
            "Risk-adjusted plan",
        ],
        size=13,
    )
    add_text(s, 0.8, 4.25, 8.0, 0.8, f"Repository: {REPO_LINK}", size=14, color="muted")
    add_footer(s, number)
    number += 1

    # 2 Executive narrative
    s = slide_base(prs)
    add_title(s, "Executive Narrative", "Simple story: pain -> proof -> decision", "Speaker A", "1:00")
    cols = [
        ("Pain", "Users repeatedly retype ingredient lists, losing time and missing items."),
        ("Proof", "Pilot data shows 41% faster list creation vs manual baseline."),
        ("Decision", "Continue direction; prioritize screenshot-noise reduction."),
    ]
    x = 0.75
    for title, body in cols:
        add_card(s, x, 1.8, 4.0, 4.7, title)
        add_text(s, x + 0.28, 2.45, 3.45, 3.6, body, size=17)
        x += 4.2
    add_footer(s, number)
    number += 1

    # 3 customer analysis
    s = slide_base(prs)
    add_title(s, "Customer Deep Analysis", "Primary fit: roommate undergrads + ELS students", "Speaker A", "1:05")
    add_card(s, 0.75, 1.75, 12.0, 4.8, "Segment Prioritization Matrix")
    headers = ["Segment", "Fit", "Reason"]
    add_text(s, 1.0, 2.1, 3.2, 0.4, headers[0], size=14, color="accent", bold=True)
    add_text(s, 4.4, 2.1, 0.9, 0.4, headers[1], size=14, color="accent", bold=True)
    add_text(s, 5.4, 2.1, 6.8, 0.4, headers[2], size=14, color="accent", bold=True)
    rows = [
        ("Roommate undergrad cooks", "5/5", "High frequency, screenshot-heavy, manual workflow pain"),
        ("ELS/international student cooks", "5/5", "Mixed-language ingredient friction + translation effort"),
        ("Recipe app power users", "2/5", "Already have structured recipe systems"),
        ("Delivery-first users", "1/5", "Low grocery-list building frequency"),
    ]
    y = 2.55
    for seg, fit, reason in rows:
        add_card(s, 0.95, y - 0.07, 11.6, 0.82)
        add_text(s, 1.15, y, 3.1, 0.5, seg, size=13)
        add_text(s, 4.45, y, 0.9, 0.5, fit, size=13, color="accent", bold=True)
        add_text(s, 5.45, y, 6.8, 0.5, reason, size=13, color="muted")
        y += 0.92
    add_footer(s, number)
    number += 1

    # 4 customer conversations
    s = slide_base(prs)
    add_title(s, "Voice of Customer", "3 real roommate interviews", "Speaker B", "1:15")
    cards = [
        ("U1 | Undergrad roommate", "1m 56s", "\"Faster than Notes, but I verify fractions.\""),
        ("U2 | Undergrad roommate", "2m 18s", "\"Great import flow. Social text noise is the pain.\""),
        ("U3 | ELS roommate from Peru", "2m 04s", "\"Checklist persistence matters most. Keep it simple.\""),
    ]
    x = 0.7
    for title, time, quote in cards:
        add_card(s, x, 1.9, 4.1, 4.7, title)
        metric(s, x + 0.25, 2.45, 1.65, 1.4, time, "Time to list")
        add_text(s, x + 2.05, 2.6, 1.7, 0.6, "Interview\ncompleted", size=12, color="accent", bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + 0.25, 4.15, 3.5, 2.1, quote, size=16)
        x += 4.24
    add_footer(s, number)
    number += 1

    # 5 hypothesis + falsification chart
    s = slide_base(prs)
    add_title(s, "Hypothesis and Falsification", "We actively tried to prove ourselves wrong", "Speaker B", "1:20")
    add_card(s, 0.7, 1.75, 5.7, 4.9, "Hypothesis")
    add_text(
        s,
        1.0,
        2.35,
        5.1,
        1.6,
        "If list creation speed improves by >=30%, users adopt weekly even with imperfect OCR.",
        size=18,
        bold=True,
    )
    add_bullets(
        s,
        1.0,
        4.1,
        5.0,
        2.2,
        [
            "Manual baseline: 3m 34s",
            "App overall: 2m 06s (41% faster)",
            "Screenshot subset: 28% faster (risk boundary)",
        ],
        size=15,
    )
    s.shapes.add_picture(str(assets["falsification"]), Inches(6.7), Inches(1.75), width=Inches(6.0), height=Inches(4.9))
    add_footer(s, number)
    number += 1

    # 6 2x2
    s = slide_base(prs)
    add_title(s, "2x2 Differentiation", "Positioning for business leaders", "Speaker B", "1:00")
    s.shapes.add_picture(str(assets["two_by_two"]), Inches(0.8), Inches(1.65), width=Inches(8.1), height=Inches(4.9))
    add_card(s, 9.2, 1.65, 3.9, 4.9, "Strategic Readout")
    add_bullets(
        s,
        9.45,
        2.25,
        3.35,
        3.8,
        [
            "Gap: flexible input + structured list output in one flow",
            "Competitors split this value across separate tools",
            "Winning condition: keep speed and trust under noisy inputs",
        ],
        size=15,
    )
    add_footer(s, number)
    number += 1

    # 7 Demo visuals
    s = slide_base(prs)
    add_title(s, "Product Demo Storyboard", "Actual workflow shown as product sequence", "Speaker C", "1:10")
    add_text(
        s,
        0.8,
        1.55,
        12.0,
        0.5,
        "Scan -> Review -> Shop (the same sequence used in testing and customer sessions)",
        size=15,
        color="muted",
    )
    s.shapes.add_picture(str(assets["demo_scan"]), Inches(0.7), Inches(1.95), width=Inches(4.0), height=Inches(4.8))
    s.shapes.add_picture(str(assets["demo_edit"]), Inches(4.75), Inches(1.95), width=Inches(4.0), height=Inches(4.8))
    s.shapes.add_picture(str(assets["demo_shop"]), Inches(8.8), Inches(1.95), width=Inches(4.0), height=Inches(4.8))
    metric(s, 0.9, 6.0, 2.1, 1.0, "0.82", "Median OCR confidence")
    metric(s, 3.2, 6.0, 2.1, 1.0, "91.3%", "Field accuracy")
    metric(s, 5.5, 6.0, 2.1, 1.0, "26.9s", "Median saved-list time")
    add_footer(s, number)
    number += 1

    # 8 System architecture
    s = slide_base(prs)
    add_title(s, "System Architecture and Leverage Points", "Where we intervene in the larger system", "Speaker B", "1:05")
    s.shapes.add_picture(str(assets["architecture"]), Inches(0.7), Inches(1.7), width=Inches(8.3), height=Inches(4.9))
    add_card(s, 9.25, 1.7, 3.8, 4.9, "Leverage Points")
    add_bullets(
        s,
        9.5,
        2.25,
        3.3,
        3.9,
        [
            "LP1 OCR quality",
            "LP2 Parser precision",
            "LP3 Edit speed",
            "LP4 Checklist persistence",
            "Target: translation bottleneck",
        ],
        size=15,
    )
    add_footer(s, number)
    number += 1

    # 9 technical process
    s = slide_base(prs)
    add_title(s, "Technical Process Evidence", "Built phase-by-phase, not one-shot", "Speaker C", "1:00")
    s.shapes.add_picture(str(assets["pipeline"]), Inches(0.8), Inches(1.65), width=Inches(12.0), height=Inches(3.35))
    add_card(s, 0.8, 5.15, 5.95, 1.45, "Process Signals Casey Cares About")
    add_bullets(
        s,
        1.05,
        5.55,
        5.4,
        0.95,
        [
            "Roadmap checklists completed by phase",
            "Docs updated as implementation evolved",
            "Core logic split for testability",
        ],
        size=14,
    )
    add_card(s, 6.9, 5.15, 5.9, 1.45, "Artifacts")
    add_bullets(
        s,
        7.15,
        5.55,
        5.3,
        0.95,
        [
            "PRD, MVP, architecture, roadmap, changelog",
            "Rubric map + artifact index",
            "Evidence docs linked to decisions",
        ],
        size=14,
    )
    add_footer(s, number)
    number += 1

    # 10 AI infra + debugging
    s = slide_base(prs)
    add_title(s, "AI Infrastructure and Debugging Discipline", "What technical leaders evaluate", "Speaker C", "0:55")
    add_card(s, 0.8, 1.7, 6.0, 4.9, "Infrastructure")
    add_bullets(
        s,
        1.05,
        2.25,
        5.4,
        3.8,
        [
            "AI docs structure (`aiDocs/`) and MCP checklist",
            "Secret-safe `.gitignore` and reproducible scripts",
            "Cross-platform constraints documented",
            "Shared `RecipeCore` for parser + fixtures",
        ],
        size=15,
    )
    add_card(s, 7.0, 1.7, 5.4, 2.35, "Structured Logs")
    add_text(
        s,
        7.25,
        2.15,
        4.9,
        1.75,
        '{"event":"ocr_low_confidence","confidence":"0.66"}\n{"event":"ocr_completed","ingredient_count":"17"}',
        size=13,
        color="muted",
    )
    add_card(s, 7.0, 4.25, 5.4, 2.35, "Test-Log-Fix")
    add_bullets(
        s,
        7.25,
        4.7,
        4.9,
        1.65,
        [
            "Reproduce -> inspect log -> patch -> rerun fixtures",
            "`scripts/run-cli-tests.ps1` + JSON fixtures",
        ],
        size=14,
    )
    add_footer(s, number)
    number += 1

    # 11 risk and gates
    s = slide_base(prs)
    add_title(s, "Success/Failure Plan and Decision Gates", "Business-governed execution", "Speaker A", "0:55")
    add_card(s, 0.8, 1.8, 5.9, 4.8, "Success Indicators")
    add_bullets(
        s,
        1.05,
        2.3,
        5.3,
        3.9,
        [
            "Median list creation remains <30s",
            "Field accuracy remains >=90%",
            "Weekly reuse intent >=50%",
            "Correction count trend decreases",
        ],
        size=15,
    )
    add_card(s, 6.95, 1.8, 5.55, 4.8, "Failure Triggers + Pivots")
    add_bullets(
        s,
        7.2,
        2.3,
        5.0,
        3.9,
        [
            "If speed edge disappears, narrow use case",
            "If trust remains low, shift to semi-manual assist",
            "If retention fails, pause expansion and re-scope",
            "Current top risk: social screenshot noise",
        ],
        size=15,
    )
    add_footer(s, number)
    number += 1

    # 12 30-day plan
    s = slide_base(prs)
    add_title(s, "30-Day Plan Before Final", "Focused execution, no scope creep", "Speaker A", "0:50")
    add_card(s, 0.8, 1.75, 12.0, 4.9, "Milestones")
    milestones = [
        ("Week 1", "Screenshot-noise parser hardening + fixture expansion"),
        ("Week 2", "Additional target-user sessions (student + ELS heavy)"),
        ("Week 3", "Measure correction trends and weekly reuse intent"),
        ("Week 4", "Demo polish + final decision checkpoint"),
    ]
    y = 2.35
    for week, desc in milestones:
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.15), Inches(y), Inches(0.25), Inches(0.25))
        dot.fill.solid()
        dot.fill.fore_color.rgb = THEME["accent"]
        dot.line.fill.background()
        add_text(s, 1.5, y - 0.02, 1.3, 0.3, week, size=14, color="accent", bold=True)
        add_text(s, 2.8, y - 0.03, 9.5, 0.5, desc, size=15)
        y += 1.02
    add_footer(s, number)
    number += 1

    # 13 final ask
    s = slide_base(prs)
    add_title(s, "Final Ask to Leadership", "Greenlight continuation with constrained scope", "Speaker A", "0:40")
    add_text(
        s,
        0.9,
        2.05,
        11.6,
        1.3,
        "We have evidence of customer value, a tested hypothesis, and disciplined technical process.\nWe are asking for approval to continue on the current narrow strategy.",
        size=22,
        bold=True,
    )
    add_card(s, 0.9, 3.85, 12.0, 2.4, "Requested Decision")
    add_bullets(
        s,
        1.2,
        4.35,
        11.0,
        1.8,
        [
            "1) Keep scope centered on translation bottleneck",
            "2) Prioritize screenshot-noise reduction over feature expansion",
            "3) Re-evaluate with retention and correction metrics at final milestone",
        ],
        size=17,
    )
    add_footer(s, number)
    number += 1

    # 14 appendix artifacts
    s = slide_base(prs)
    add_title(s, "Appendix: Required Artifact Submission", "Everything requested in rubric is in-repo", "Backup", "")
    add_card(s, 0.8, 1.75, 12.0, 4.9, "Artifact Checklist")
    add_bullets(
        s,
        1.05,
        2.3,
        11.2,
        3.9,
        [
            f"GitHub repository link: {REPO_LINK}",
            "Deep customer analysis",
            "Founding hypothesis",
            "3 real customer conversation documents",
            "1 falsification test report",
            "1 2x2 differentiation grid",
            "System architecture with leverage points and target layer",
            "Technical process evidence, logging/debugging, and roadmap progression",
        ],
        size=15,
    )
    add_footer(s, number)
    number += 1

    # 15 artifact map
    s = slide_base(prs)
    add_title(s, "Reference Map", "Where reviewers can verify every claim", "Backup", "")
    add_card(s, 0.8, 1.75, 12.0, 4.9, "Key Files")
    add_bullets(
        s,
        1.05,
        2.2,
        11.2,
        3.9,
        [
            "aiDocs/founding-hypothesis.md",
            "aiDocs/falsification-test.md",
            "aiDocs/differentiation-2x2.md",
            "aiDocs/system-architecture-leverage.md",
            "aiDocs/evidence/customer-conversation-u1.md / u2.md / u3.md",
            "aiDocs/debugging.md and scripts/run-cli-tests.ps1",
        ],
        size=15,
    )
    add_footer(s, number)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)


def main() -> None:
    assets = ensure_assets()
    build_presentation(assets)
    print(f"Created polished deck: {OUTPUT}")
    print(f"Created visual assets in: {ASSET_DIR}")


if __name__ == "__main__":
    main()
